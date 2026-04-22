from __future__ import annotations

import io
import os
import re
import base64
import logging
from datetime import date
from flask import Flask, send_file, request, send_from_directory, jsonify
from werkzeug.exceptions import RequestEntityTooLarge
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn

# ── Speed: fast-compress on save ──────────────────────────────────────────────
# A .pptx file is internally a compressed archive (Microsoft's Open XML
# format). python-pptx defaults to the slowest/tightest compression setting,
# which spends ~180 ms per deck purely on compression math. Dropping from
# level 6 to level 1 saves ~40 % of that save time for only ~1 % more bytes —
# the output is still a normal, fully-compressed .pptx that opens in
# PowerPoint exactly like any other deck. We swap python-pptx's internal
# writer rather than vendor / subclass the library.
import zipfile as _zipfile
from pptx.opc.serialized import _ZipPkgWriter  # noqa: E402
from pptx.util import lazyproperty as _lazyproperty  # noqa: E402


@_lazyproperty
def _fast_zipf(self):
    return _zipfile.ZipFile(
        self._pkg_file, "w",
        compression=_zipfile.ZIP_DEFLATED, compresslevel=1,
        strict_timestamps=False,
    )


_ZipPkgWriter._zipf = _fast_zipf

app = Flask(__name__)
# Allow generous room for the rep's base64-encoded photo in form data.
# Default Werkzeug caps are ~500 KB which rejects most phone-camera photos.
app.config["MAX_CONTENT_LENGTH"]    = 25 * 1024 * 1024  # 25 MB request body
app.config["MAX_FORM_MEMORY_SIZE"]  = 25 * 1024 * 1024  # 25 MB form fields

# ── Config ────────────────────────────────────────────────────────────────────
BASE_DIR              = os.path.dirname(os.path.abspath(__file__))
TEMPLATE_2TERM        = os.path.join(BASE_DIR, "Datasite Proposal_6_9_options _small.pptx")
TEMPLATE_1TERM        = os.path.join(BASE_DIR, "3_Mos_Pricing_Diligence_small.pptx")
TEMPLATE_1TERM_TIERED = os.path.join(BASE_DIR, "tiered_1term_small.pptx")
TEMPLATE_2TERM_TIERED = os.path.join(BASE_DIR, "tiered_2term_small.pptx")
TEMPLATE_3TERM_TIERED = os.path.join(BASE_DIR, "tiered_3term_small.pptx")

PORT = int(os.environ.get("PORT", "5050"))

# ── Industry slides ───────────────────────────────────────────────────────────
# Each entry is `display_name -> slide index inside every template`.
# To add a new industry:
#   1. Open each of the 4 PPTX templates in PowerPoint.
#   2. Duplicate the Technology or Healthcare slide (Cmd-D) and edit its
#      copy/images for the new industry. Place it AFTER the existing
#      industry slides so existing indexes don't shift.
#   3. Add a single line below mapping the industry name to its new slide
#      index, then restart the server. The form will pick it up automatically.
INDUSTRY_SLIDES: dict[str, int] = {
    "Technology":      6,
    "Healthcare":      7,
    "Real Estate":     8,
    "Financial":       9,
    "Industrials":    10,
    "Energy":         11,
    "Consumer Retail": 12,
}

# Headers the SSO reverse proxy may forward with the authenticated user.
# Safe to leave as-is; falls back to "-" if nothing is present.
SSO_USER_HEADERS = ("X-Forwarded-User", "X-Remote-User", "X-Auth-Request-Email", "X-User-Email")

# ── Logging ───────────────────────────────────────────────────────────────────
logging.basicConfig(
    level=logging.INFO,
    format="%(asctime)s  %(levelname)s  %(message)s",
    datefmt="%Y-%m-%d %H:%M:%S",
)
log = logging.getLogger("proposal-builder")


def current_user():
    for h in SSO_USER_HEADERS:
        v = request.headers.get(h)
        if v:
            return v.strip()
    return "-"


# ── Startup check ─────────────────────────────────────────────────────────────
TEMPLATES_REQUIRED = {
    "1-term":        TEMPLATE_1TERM,
    "2-term":        TEMPLATE_2TERM,
    "1-term tiered": TEMPLATE_1TERM_TIERED,
    "2-term tiered": TEMPLATE_2TERM_TIERED,
    "3-term tiered": TEMPLATE_3TERM_TIERED,
}
log.info("Base dir: %s", BASE_DIR)
for label, path in TEMPLATES_REQUIRED.items():
    log.info("Template %-15s %s  %s", label, "OK " if os.path.exists(path) else "MISSING", path)

MINIMUMS = {3: 1500, 4: 2000, 6: 3000, 9: 4500, 12: 6000}

# Matches any M/D/YY or M/D/YYYY token, bounded by non-digits on each side.
# Uses a capture group so we can re-insert surrounding context untouched.
DATE_RE = re.compile(r"(?<!\d)(\d{1,2}/\d{1,2}/\d{2,4})(?!\d)")


def minimum(term: int) -> int:
    return MINIMUMS.get(int(term), 500 * int(term))


# ── Safe form coercion ────────────────────────────────────────────────────────
class FormError(ValueError):
    pass


def form_int(name: str, default: int, min_value: int | None = None,
             max_value: int | None = None) -> int:
    raw = (request.form.get(name, "") or "").strip()
    if raw == "":
        return default
    try:
        v = int(float(raw))   # accept "3" or "3.0"
    except ValueError:
        raise FormError(f"Field '{name}' must be a whole number (got {raw!r}).")
    if min_value is not None and v < min_value:
        raise FormError(f"Field '{name}' must be ≥ {min_value} (got {v}).")
    if max_value is not None and v > max_value:
        raise FormError(f"Field '{name}' must be ≤ {max_value} (got {v}).")
    return v


def form_float(name: str, default: float, min_value: float | None = None) -> float:
    raw = (request.form.get(name, "") or "").strip()
    if raw == "":
        return default
    try:
        v = float(raw)
    except ValueError:
        raise FormError(f"Field '{name}' must be a number (got {raw!r}).")
    if min_value is not None and v < min_value:
        raise FormError(f"Field '{name}' must be ≥ {min_value} (got {v}).")
    return v


# ── PPTX helpers ──────────────────────────────────────────────────────────────
def set_cell(cell, text: str) -> None:
    """Replace the visible text of a table cell, preserving its first run's formatting."""
    for para in cell.text_frame.paragraphs:
        for i, run in enumerate(para.runs):
            run.text = text if i == 0 else ""
        if para.runs:
            return
    for t in cell._tc.findall('.//' + qn('a:t')):
        t.text = text
        return


def _set_paragraph_text(paragraph, text: str) -> None:
    """Replace paragraph text by writing into its first run (preserves font,
    size, color, bold, etc.) and blanking any trailing runs."""
    runs = list(paragraph.runs)
    if runs:
        runs[0].text = text
        for r in runs[1:]:
            r.text = ""
    elif text:
        paragraph.add_run().text = text


def set_rep_caption(shape, name: str, title: str) -> None:
    """Populate the rep-card caption the same way the template was designed:
    paragraph 0 = name (large), paragraph 1 = title (small). Any extra
    paragraphs in the template (e.g. city line) are blanked so we don't leak
    the original designer's info. If the caller passes empty strings, the
    corresponding line is blanked rather than left untouched."""
    if not shape.has_text_frame:
        return
    paragraphs = list(shape.text_frame.paragraphs)
    if len(paragraphs) >= 1:
        _set_paragraph_text(paragraphs[0], name)
    if len(paragraphs) >= 2:
        _set_paragraph_text(paragraphs[1], title)
    for p in paragraphs[2:]:
        for r in p.runs:
            r.text = ""


def swap_picture_image(pic_shape, host_slide, new_image_bytes: bytes) -> None:
    """Hot-swap the image bytes of an existing PICTURE shape, leaving its
    position, size, and any container (group) intact. Works by adding the new
    image as a new image part on the slide and re-pointing the picture's
    blip reference at it; the staging shape used to register the part is
    deleted immediately so it never renders."""
    staging = host_slide.shapes.add_picture(
        io.BytesIO(new_image_bytes), 0, 0, Inches(0.01), Inches(0.01)
    )
    blip_ns = qn("a:blip")
    embed_attr = qn("r:embed")
    new_rId = staging._element.blipFill.find(blip_ns).get(embed_attr)
    staging._element.getparent().remove(staging._element)
    pic_shape._element.blipFill.find(blip_ns).set(embed_attr, new_rId)


def remove_slide(prs, index: int) -> None:
    xml_slides = prs.slides._sldIdLst
    slide = prs.slides[index]
    rId = None
    for rel in prs.part.rels.values():
        if rel._target == slide.part:
            rId = rel.rId
            break
    if rId:
        prs.part.drop_rel(rId)
    del xml_slides[index]


def keep_only_industry_slide(prs, industry: str | None) -> None:
    """Strip every industry slide except the one matching `industry`. Indexes
    are read from INDUSTRY_SLIDES so adding a new industry never requires
    touching this function. Slides are removed highest-index-first so earlier
    indexes don't shift mid-loop."""
    keep_idx = INDUSTRY_SLIDES.get(industry) if industry else None
    drop = sorted(
        (idx for name, idx in INDUSTRY_SLIDES.items() if idx != keep_idx),
        reverse=True,
    )
    for idx in drop:
        remove_slide(prs, idx)


def tiered_cost(pages: int, r1: float, r2: float, r3: float = 0.0,
                upper1: int = 9999, upper2: int | None = 19999) -> float:
    """Flat-rate-at-volume: total pages × the rate for that volume band.

    3-tier (default):  up to upper1 → r1; (upper1+1)–upper2 → r2; (upper2+1)+ → r3.
    2-tier:            pass upper2=None — up to upper1 → r1; (upper1+1)+ → r2.
    """
    if upper2 is None:
        return pages * (r1 if pages <= upper1 else r2)
    if pages <= upper1:
        return pages * r1
    if pages <= upper2:
        return pages * r2
    return pages * r3


def tier_band_labels(upper1: int, upper2: int | None = None) -> tuple[str, ...]:
    """Slide labels that always match the boundaries actually used.
    Returns 2 labels for a 2-tier deal, 3 for a 3-tier deal."""
    if upper2 is None:
        return (
            f"Up to {upper1:,}",
            f"{upper1 + 1:,}+",
        )
    return (
        f"Up to {upper1:,}",
        f"{upper1 + 1:,}–{upper2:,}",
        f"{upper2 + 1:,}+",
    )


def remove_table_row(table, row_index: int) -> None:
    """Drop a single row from a python-pptx table by removing its <a:tr>
    element from the underlying XML."""
    tr = list(table._tbl.tr_lst)[row_index]
    table._tbl.remove(tr)


# ── Deck builder ──────────────────────────────────────────────────────────────
def build_deck(num_terms, term_rates, prep, pages1, pages2, first_names, company,
               industry=None, pricing_type="term1",
               mtm_rate=0.15, tiered_months=3,
               num_tiers=3,
               tier1_rate=0.37, tier2_rate=0.32, tier3_rate=0.29,
               tier_upper1=9999, tier_upper2=19999,
               t2t_a_months=3, t2t_a_r1=0.37, t2t_a_r2=0.32, t2t_a_r3=0.29,
               t2t_b_months=6, t2t_b_r1=0.35, t2t_b_r2=0.30, t2t_b_r3=0.27,
               t3t_a_months=3, t3t_a_r1=0.59, t3t_a_r2=0.59, t3t_a_r3=0.59,
               t3t_b_months=6, t3t_b_r1=0.71, t3t_b_r2=0.71, t3t_b_r3=0.71,
               t3t_c_months=9, t3t_c_r1=0.71, t3t_c_r2=0.71, t3t_c_r3=0.71,
               rep_photo_bytes=None, signer_name="", signer_title=""):
    terms       = sorted(term_rates.keys())
    prep        = int(prep)
    pages1      = int(pages1) if pages1 else 5000
    pages2      = int(pages2) if pages2 else 10000
    first_names = (first_names or "").strip() or "Team"
    company     = (company or "").strip()
    today_short = date.today().strftime("%-m/%-d/%y")

    if pricing_type == "tiered":
        template = TEMPLATE_1TERM_TIERED
    elif pricing_type == "tiered2":
        template = TEMPLATE_2TERM_TIERED
    elif pricing_type == "tiered3":
        template = TEMPLATE_3TERM_TIERED
    elif num_terms == 2:
        template = TEMPLATE_2TERM
    else:
        template = TEMPLATE_1TERM
    prs = Presentation(template)

    keep_only_industry_slide(prs, industry)
    pricing_slide = prs.slides[-1]

    # ── Slide 1: Cover ────────────────────────────────────────────────────────
    cover = prs.slides[0]
    for shape in cover.shapes:
        if shape.has_text_frame:
            paras = shape.text_frame.paragraphs
            if len(paras) >= 3 and 'Prepared for' in paras[0].text:
                if paras[1].runs:
                    paras[1].runs[0].text = company or first_names
                    for r in paras[1].runs[1:]:
                        r.text = ""
                if paras[2].runs:
                    paras[2].runs[0].text = date.today().strftime("%B %d, %Y")
                    for r in paras[2].runs[1:]:
                        r.text = ""
                break

    # ── Slide 2: Dear …, rep photo, signer block ─────────────────────────────
    slide2 = prs.slides[1]
    for shape in slide2.shapes:
        if shape.name == "TextBox 89":
            para = shape.text_frame.paragraphs[0]
            if para.runs:
                para.runs[0].text = f"Dear {first_names},"
                for run in para.runs[1:]:
                    run.text = ""
            break

    # Replace the hardcoded "rep card" baked into slide 2 (a Group containing
    # a headshot picture + a caption box). This is the prominent right-side
    # photo every rep sees — without this swap it would always show the
    # template designer's headshot and name.
    if rep_photo_bytes or signer_name or signer_title:
        for shape in slide2.shapes:
            if shape.shape_type != MSO_SHAPE_TYPE.GROUP:
                continue
            for sub in shape.shapes:
                if rep_photo_bytes and sub.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    try:
                        swap_picture_image(sub, slide2, rep_photo_bytes)
                    except Exception as e:
                        log.warning("Rep photo swap failed: %s", e)
                elif (signer_name or signer_title) and sub.has_text_frame:
                    try:
                        set_rep_caption(sub, signer_name, signer_title)
                    except Exception as e:
                        log.warning("Rep caption update failed: %s", e)

    # ── Pricing slide ─────────────────────────────────────────────────────────
    for shape in pricing_slide.shapes:

        if shape.name == "Text Placeholder 1":
            runs = shape.text_frame.paragraphs[0].runs
            if runs:
                if pricing_type == "mtm":
                    runs[0].text = "Pricing option: Warehouse"
                elif pricing_type == "tiered":
                    runs[0].text = f"Pricing option: {tiered_months}-month"
                elif pricing_type == "tiered2":
                    runs[0].text = f"Pricing options: {t2t_a_months}/{t2t_b_months} Months"
                elif pricing_type == "tiered3":
                    runs[0].text = f"Pricing options: {t3t_a_months}/{t3t_b_months}/{t3t_c_months} Months"
                elif num_terms == 1:
                    runs[0].text = f"Pricing option: {terms[0]}-month"
                else:
                    runs[0].text = f"Pricing options: {'/'.join(str(t) for t in terms)} Months"
                for r in runs[1:]:
                    r.text = ""

        elif shape.name == "Table 1":
            t = shape.table
            set_cell(t.cell(0, 1), f"{prep} days")
            if pricing_type == "mtm":
                set_cell(t.cell(2, 2), "Warehouse")
                set_cell(t.cell(3, 1), "All pages")
                set_cell(t.cell(3, 2), f"${mtm_rate:.2f}/page")
                # Row 4 in the 1-term template is the Project-management row —
                # leave it untouched.
            elif pricing_type == "tiered":
                upper2_eff = tier_upper2 if num_tiers == 3 else None
                labels = tier_band_labels(tier_upper1, upper2_eff)
                set_cell(t.cell(2, 2), f"{tiered_months} Months")
                set_cell(t.cell(3, 1), labels[0])
                set_cell(t.cell(3, 2), f"${tier1_rate:.2f}/page")
                set_cell(t.cell(4, 1), labels[1])
                set_cell(t.cell(4, 2), f"${tier2_rate:.2f}/page")
                if num_tiers == 3:
                    set_cell(t.cell(5, 1), labels[2])
                    set_cell(t.cell(5, 2), f"${tier3_rate:.2f}/page")
                else:
                    remove_table_row(t, 5)
            elif pricing_type == "tiered2":
                upper2_eff = tier_upper2 if num_tiers == 3 else None
                labels = tier_band_labels(tier_upper1, upper2_eff)
                set_cell(t.cell(2, 2), f"{t2t_a_months} Months")
                set_cell(t.cell(2, 3), f"{t2t_b_months} Months")
                set_cell(t.cell(3, 1), labels[0])
                set_cell(t.cell(3, 2), f"${t2t_a_r1:.2f}/page")
                set_cell(t.cell(3, 3), f"${t2t_b_r1:.2f}/page")
                set_cell(t.cell(4, 1), labels[1])
                set_cell(t.cell(4, 2), f"${t2t_a_r2:.2f}/page")
                set_cell(t.cell(4, 3), f"${t2t_b_r2:.2f}/page")
                if num_tiers == 3:
                    set_cell(t.cell(5, 1), labels[2])
                    set_cell(t.cell(5, 2), f"${t2t_a_r3:.2f}/page")
                    set_cell(t.cell(5, 3), f"${t2t_b_r3:.2f}/page")
                else:
                    remove_table_row(t, 5)
            elif pricing_type == "tiered3":
                upper2_eff = tier_upper2 if num_tiers == 3 else None
                labels = tier_band_labels(tier_upper1, upper2_eff)
                set_cell(t.cell(2, 2), f"{t3t_a_months} Months")
                set_cell(t.cell(2, 3), f"{t3t_b_months} Months")
                set_cell(t.cell(2, 4), f"{t3t_c_months} Months")
                set_cell(t.cell(3, 1), labels[0])
                set_cell(t.cell(3, 2), f"${t3t_a_r1:.2f}/page")
                set_cell(t.cell(3, 3), f"${t3t_b_r1:.2f}/page")
                set_cell(t.cell(3, 4), f"${t3t_c_r1:.2f}/page")
                set_cell(t.cell(4, 1), labels[1])
                set_cell(t.cell(4, 2), f"${t3t_a_r2:.2f}/page")
                set_cell(t.cell(4, 3), f"${t3t_b_r2:.2f}/page")
                set_cell(t.cell(4, 4), f"${t3t_c_r2:.2f}/page")
                if num_tiers == 3:
                    set_cell(t.cell(5, 1), labels[2])
                    set_cell(t.cell(5, 2), f"${t3t_a_r3:.2f}/page")
                    set_cell(t.cell(5, 3), f"${t3t_b_r3:.2f}/page")
                    set_cell(t.cell(5, 4), f"${t3t_c_r3:.2f}/page")
                else:
                    remove_table_row(t, 5)
            elif num_terms == 1:
                term = terms[0]
                set_cell(t.cell(2, 2), f"Term {term} months")
                set_cell(t.cell(3, 1), "All pages")
                set_cell(t.cell(3, 2), f"${term_rates[term]:.2f}/page")
            else:
                for i in range(2):
                    col = i + 2
                    if i < len(terms):
                        set_cell(t.cell(2, col), f"{terms[i]} Months")
                        set_cell(t.cell(3, col), f"${term_rates[terms[i]]:.2f}/page")
                    else:
                        set_cell(t.cell(2, col), "")
                        set_cell(t.cell(3, col), "")

        elif shape.name == "Table 5":
            t = shape.table
            if pricing_type == "mtm":
                set_cell(t.cell(0, 0), "Warehouse cost estimates")
                set_cell(t.cell(0, 1), "Costs")
                set_cell(t.cell(1, 0), f"Cost assuming {pages1:,} pages")
                set_cell(t.cell(1, 1), f"${pages1 * mtm_rate:,.0f}")
                set_cell(t.cell(2, 0), f"Cost assuming {pages2:,} pages")
                set_cell(t.cell(2, 1), f"${pages2 * mtm_rate:,.0f}")

            elif pricing_type == "tiered":
                upper2_eff = tier_upper2 if num_tiers == 3 else None
                set_cell(t.cell(0, 0), "Cost estimates")
                set_cell(t.cell(0, 1), "Costs")
                set_cell(t.cell(1, 0), f"Cost assuming {pages1:,} pages")
                set_cell(t.cell(1, 1), f"${tiered_cost(pages1, tier1_rate, tier2_rate, tier3_rate, tier_upper1, upper2_eff):,.0f}")
                set_cell(t.cell(2, 0), f"Cost assuming {pages2:,} pages")
                set_cell(t.cell(2, 1), f"${tiered_cost(pages2, tier1_rate, tier2_rate, tier3_rate, tier_upper1, upper2_eff):,.0f}")

            elif pricing_type == "tiered2":
                upper2_eff = tier_upper2 if num_tiers == 3 else None
                set_cell(t.cell(0, 0), "Cost estimates")
                set_cell(t.cell(0, 1), f"{t2t_a_months} Month")
                set_cell(t.cell(0, 2), f"{t2t_b_months} Month")
                set_cell(t.cell(1, 0), f"Cost assuming {pages1:,} pages")
                set_cell(t.cell(1, 1), f"${tiered_cost(pages1, t2t_a_r1, t2t_a_r2, t2t_a_r3, tier_upper1, upper2_eff):,.0f}")
                set_cell(t.cell(1, 2), f"${tiered_cost(pages1, t2t_b_r1, t2t_b_r2, t2t_b_r3, tier_upper1, upper2_eff):,.0f}")
                set_cell(t.cell(2, 0), f"Cost assuming {pages2:,} pages")
                set_cell(t.cell(2, 1), f"${tiered_cost(pages2, t2t_a_r1, t2t_a_r2, t2t_a_r3, tier_upper1, upper2_eff):,.0f}")
                set_cell(t.cell(2, 2), f"${tiered_cost(pages2, t2t_b_r1, t2t_b_r2, t2t_b_r3, tier_upper1, upper2_eff):,.0f}")

            elif pricing_type == "tiered3":
                upper2_eff = tier_upper2 if num_tiers == 3 else None
                set_cell(t.cell(0, 0), "Cost estimate")
                set_cell(t.cell(0, 1), f"{t3t_a_months} Month")
                set_cell(t.cell(0, 2), f"{t3t_b_months} Month")
                set_cell(t.cell(0, 3), f"{t3t_c_months} Month")
                set_cell(t.cell(1, 0), f"Cost assuming {pages1:,} pages")
                set_cell(t.cell(1, 1), f"${tiered_cost(pages1, t3t_a_r1, t3t_a_r2, t3t_a_r3, tier_upper1, upper2_eff):,.0f}")
                set_cell(t.cell(1, 2), f"${tiered_cost(pages1, t3t_b_r1, t3t_b_r2, t3t_b_r3, tier_upper1, upper2_eff):,.0f}")
                set_cell(t.cell(1, 3), f"${tiered_cost(pages1, t3t_c_r1, t3t_c_r2, t3t_c_r3, tier_upper1, upper2_eff):,.0f}")
                set_cell(t.cell(2, 0), f"Cost assuming {pages2:,} pages")
                set_cell(t.cell(2, 1), f"${tiered_cost(pages2, t3t_a_r1, t3t_a_r2, t3t_a_r3, tier_upper1, upper2_eff):,.0f}")
                set_cell(t.cell(2, 2), f"${tiered_cost(pages2, t3t_b_r1, t3t_b_r2, t3t_b_r3, tier_upper1, upper2_eff):,.0f}")
                set_cell(t.cell(2, 3), f"${tiered_cost(pages2, t3t_c_r1, t3t_c_r2, t3t_c_r3, tier_upper1, upper2_eff):,.0f}")

            elif num_terms == 1:
                term     = terms[0]
                rate     = term_rates[term]
                set_cell(t.cell(0, 0), "Term 1 cost estimates")
                set_cell(t.cell(0, 1), "Costs")
                cost1_raw = pages1 * rate
                cost1     = max(cost1_raw, minimum(term))
                flag1     = " *" if cost1_raw < minimum(term) else ""
                set_cell(t.cell(1, 0), f"Cost assuming {pages1:,} pages")
                set_cell(t.cell(1, 1), f"${cost1:,.0f}{flag1}")
                cost2_raw = pages2 * rate
                cost2     = max(cost2_raw, minimum(term))
                flag2     = " *" if cost2_raw < minimum(term) else ""
                set_cell(t.cell(2, 0), f"Cost assuming {pages2:,} pages")
                set_cell(t.cell(2, 1), f"${cost2:,.0f}{flag2}")

            else:
                set_cell(t.cell(0, 0), "Cost estimate")
                for i, term in enumerate(terms[:2]):
                    set_cell(t.cell(0, i + 1), f"{term} Month")
                set_cell(t.cell(1, 0), f"Cost assuming {pages1:,} pages")
                for i, term in enumerate(terms[:2]):
                    rate     = term_rates[term]
                    cost_raw = pages1 * rate
                    final    = max(cost_raw, minimum(term))
                    flag     = " *" if cost_raw < minimum(term) else ""
                    set_cell(t.cell(1, i + 1), f"${final:,.0f}{flag}")
                set_cell(t.cell(2, 0), f"Cost assuming {pages2:,} pages")
                for i, term in enumerate(terms[:2]):
                    rate     = term_rates[term]
                    cost_raw = pages2 * rate
                    final    = max(cost_raw, minimum(term))
                    flag     = " *" if cost_raw < minimum(term) else ""
                    set_cell(t.cell(2, i + 1), f"${final:,.0f}{flag}")

    # ── Footer dates ──────────────────────────────────────────────────────────
    # Rewrite every short date token (e.g. 4/6/26, 4/19/26) to today. This
    # replaces the old hard-coded list, which silently broke every time the
    # templates were refreshed with a newer date.
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if run.text and "/" in run.text and DATE_RE.search(run.text):
                            # Guard against "24/7/365" false positive.
                            run.text = DATE_RE.sub(
                                lambda m: today_short if m.group(1) != "24/7/365" else m.group(1),
                                run.text,
                            )

    # ── Serialize ─────────────────────────────────────────────────────────────
    slug     = (company or first_names).replace(" ", "_").replace(",", "").replace("/", "")
    filename = f"{slug}_{date.today().strftime('%m%d%Y')}.pptx"
    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf, filename


# ── Routes ────────────────────────────────────────────────────────────────────
@app.route("/")
def index():
    return send_from_directory(BASE_DIR, "index.html")


@app.route("/healthz")
def healthz():
    missing = [label for label, path in TEMPLATES_REQUIRED.items() if not os.path.exists(path)]
    if missing:
        return jsonify({"ok": False, "missing_templates": missing}), 500
    return jsonify({"ok": True})


@app.route("/industries")
def industries():
    """Render-time list of selectable industries. The form fetches this on
    page load and renders one pill per industry, so adding a new industry
    only requires editing INDUSTRY_SLIDES (no HTML changes)."""
    return jsonify({"industries": list(INDUSTRY_SLIDES.keys())})


@app.route("/generate", methods=["POST"])
def generate():
    try:
        pricing_type   = (request.form.get("pricing_type", "term1") or "term1").strip()
        num_terms      = form_int("num_terms",    1, min_value=1)
        prep           = form_int("prep",        30, min_value=0)
        pages1         = form_int("pages1",    5000, min_value=1)
        pages2         = form_int("pages2",   10000, min_value=1)
        first_names    = (request.form.get("first_names", "") or "").strip()
        company        = (request.form.get("company", "") or "").strip()
        industry       = (request.form.get("industry", "") or "").strip()

        mtm_rate       = form_float("mtm_rate",       0.15, min_value=0.0)
        tiered_months  = form_int(  "tiered_months",  3,    min_value=1)
        tier1_rate     = form_float("tier1_rate",     0.37, min_value=0.0)
        tier2_rate     = form_float("tier2_rate",     0.32, min_value=0.0)
        tier3_rate     = form_float("tier3_rate",     0.29, min_value=0.0)
        t2t_a_months   = form_int(  "t2t_a_months",   3,    min_value=1)
        t2t_a_r1       = form_float("t2t_a_r1",       0.37, min_value=0.0)
        t2t_a_r2       = form_float("t2t_a_r2",       0.32, min_value=0.0)
        t2t_a_r3       = form_float("t2t_a_r3",       0.29, min_value=0.0)
        t2t_b_months   = form_int(  "t2t_b_months",   6,    min_value=1)
        t2t_b_r1       = form_float("t2t_b_r1",       0.35, min_value=0.0)
        t2t_b_r2       = form_float("t2t_b_r2",       0.30, min_value=0.0)
        t2t_b_r3       = form_float("t2t_b_r3",       0.27, min_value=0.0)
        t3t_a_months   = form_int(  "t3t_a_months",   3,    min_value=1)
        t3t_a_r1       = form_float("t3t_a_r1",       0.59, min_value=0.0)
        t3t_a_r2       = form_float("t3t_a_r2",       0.59, min_value=0.0)
        t3t_a_r3       = form_float("t3t_a_r3",       0.59, min_value=0.0)
        t3t_b_months   = form_int(  "t3t_b_months",   6,    min_value=1)
        t3t_b_r1       = form_float("t3t_b_r1",       0.71, min_value=0.0)
        t3t_b_r2       = form_float("t3t_b_r2",       0.71, min_value=0.0)
        t3t_b_r3       = form_float("t3t_b_r3",       0.71, min_value=0.0)
        t3t_c_months   = form_int(  "t3t_c_months",   9,    min_value=1)
        t3t_c_r1       = form_float("t3t_c_r1",       0.71, min_value=0.0)
        t3t_c_r2       = form_float("t3t_c_r2",       0.71, min_value=0.0)
        t3t_c_r3       = form_float("t3t_c_r3",       0.71, min_value=0.0)

        # How many tiers? Most deals are 3-tier; some are 2-tier.
        num_tiers = form_int("num_tiers", 3, min_value=2, max_value=3)

        # Tier boundaries (shared across both tiered pricing types). Defaults
        # match the most common deal: up to 9,999 / 10,000–19,999 / 20,000+.
        # For a 2-tier deal only tier_upper1 is used.
        tier_upper1 = form_int("tier_upper1",  9999, min_value=1)
        tier_upper2 = form_int("tier_upper2", 19999, min_value=1)
        if num_tiers == 3 and tier_upper2 <= tier_upper1:
            raise FormError(
                f"Tier 2 upper bound ({tier_upper2:,}) must be larger than "
                f"Tier 1 upper bound ({tier_upper1:,})."
            )

        term_rates: dict[int, float] = {}
        if pricing_type in ("term1", "mtm"):
            m = form_int(  "term1_months", 3,    min_value=1)
            r = form_float("term1_rate",   0.37, min_value=0.0)
            term_rates[m] = r
        elif pricing_type == "term2":
            m1 = form_int(  "term2a_months", 6,    min_value=1)
            r1 = form_float("term2a_rate",   0.59, min_value=0.0)
            m2 = form_int(  "term2b_months", 9,    min_value=1)
            r2 = form_float("term2b_rate",   0.71, min_value=0.0)
            term_rates[m1] = r1
            term_rates[m2] = r2

        rep_photo_bytes = None
        rep_photo_b64 = (request.form.get("rep_photo_b64", "") or "").strip()
        if rep_photo_b64:
            try:
                payload = rep_photo_b64.split(",", 1)[1] if "," in rep_photo_b64 else rep_photo_b64
                rep_photo_bytes = base64.b64decode(payload)
            except Exception as e:
                log.warning("Rep photo decode failed: %s", e)

        signer_name  = (request.form.get("signer_name", "")  or "").strip()
        signer_title = (request.form.get("signer_title", "") or "").strip()

        buf, filename = build_deck(
            num_terms, term_rates, prep, pages1, pages2,
            first_names, company, industry,
            pricing_type=pricing_type,
            mtm_rate=mtm_rate,
            tiered_months=tiered_months,
            num_tiers=num_tiers,
            tier1_rate=tier1_rate, tier2_rate=tier2_rate, tier3_rate=tier3_rate,
            tier_upper1=tier_upper1, tier_upper2=tier_upper2,
            t2t_a_months=t2t_a_months, t2t_a_r1=t2t_a_r1, t2t_a_r2=t2t_a_r2, t2t_a_r3=t2t_a_r3,
            t2t_b_months=t2t_b_months, t2t_b_r1=t2t_b_r1, t2t_b_r2=t2t_b_r2, t2t_b_r3=t2t_b_r3,
            t3t_a_months=t3t_a_months, t3t_a_r1=t3t_a_r1, t3t_a_r2=t3t_a_r2, t3t_a_r3=t3t_a_r3,
            t3t_b_months=t3t_b_months, t3t_b_r1=t3t_b_r1, t3t_b_r2=t3t_b_r2, t3t_b_r3=t3t_b_r3,
            t3t_c_months=t3t_c_months, t3t_c_r1=t3t_c_r1, t3t_c_r2=t3t_c_r2, t3t_c_r3=t3t_c_r3,
            rep_photo_bytes=rep_photo_bytes,
            signer_name=signer_name, signer_title=signer_title,
        )
        log.info("GENERATE user=%s company=%r type=%s file=%s",
                 current_user(), company, pricing_type, filename)

        return send_file(
            buf,
            as_attachment=True,
            download_name=filename,
            mimetype="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        )

    except FormError as e:
        log.info("BAD_REQUEST user=%s err=%s", current_user(), e)
        return (
            "<!doctype html><meta charset='utf-8'>"
            "<title>Check your inputs</title>"
            "<div style='font-family:system-ui;max-width:560px;margin:60px auto;color:#2b2f36'>"
            "<h2 style='margin-bottom:8px'>Check your inputs</h2>"
            f"<p>{e}</p>"
            "<p><a href='javascript:history.back()'>← Go back</a></p>"
            "</div>",
            400,
        )
    except Exception as e:
        log.exception("GENERATE failed user=%s err=%s", current_user(), e)
        return (
            "<!doctype html><meta charset='utf-8'>"
            "<title>Something went wrong</title>"
            "<div style='font-family:system-ui;max-width:560px;margin:60px auto;color:#2b2f36'>"
            "<h2 style='margin-bottom:8px'>Something went wrong generating the deck.</h2>"
            "<p>The error has been logged. Please try again, and if it keeps failing "
            "share the time of this message with whoever supports the app.</p>"
            "<p><a href='javascript:history.back()'>← Go back</a></p>"
            "</div>",
            500,
        )


@app.errorhandler(RequestEntityTooLarge)
def too_large(_e):
    log.info("TOO_LARGE user=%s", current_user())
    return (
        "<!doctype html><meta charset='utf-8'>"
        "<title>Photo too big</title>"
        "<div style='font-family:system-ui;max-width:560px;margin:60px auto;color:#2b2f36'>"
        "<h2 style='margin-bottom:8px'>Your photo is too large to upload.</h2>"
        "<p>Try a smaller image (under ~15&nbsp;MB) or take a fresh photo at "
        "a lower resolution. Most phone cameras let you choose a smaller size.</p>"
        "<p><a href='javascript:history.back()'>← Go back</a></p>"
        "</div>",
        413,
    )


if __name__ == "__main__":
    log.info("Datasite Proposal Builder starting on http://localhost:%d", PORT)
    log.info("Press Ctrl+C to stop")
    app.run(host="0.0.0.0", port=PORT, debug=False)
