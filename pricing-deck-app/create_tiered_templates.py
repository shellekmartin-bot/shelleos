"""
Run once: creates tiered versions of both templates with 3 pricing rows.
Output: tiered_1term_small.pptx, tiered_2term_small.pptx
"""
import copy, os
from pptx import Presentation

BASE = os.path.dirname(os.path.abspath(__file__))
NS   = "http://schemas.openxmlformats.org/drawingml/2006/main"

PRICING_ROW_IDX = 3      # index of the first (and only) pricing row in both templates
PRICING_ROW_H   = 210000 # EMU per pricing row in tiered version (~0.23 inches, fits text)

def get_shape(slide, name):
    for s in slide.shapes:
        if s.name == name:
            return s
    return None

def clone_row_after(tbl, idx):
    """Deep-copy row at idx, clear its text, insert immediately after idx."""
    rows  = tbl.findall(f"{{{NS}}}tr")
    src   = rows[idx]
    new   = copy.deepcopy(src)
    for t in new.findall(f".//{{{NS}}}t"):
        t.text = ""
    src.addnext(new)

def make_tiered_template(src_file, dst_file):
    prs   = Presentation(src_file)
    slide = prs.slides[-1]

    tbl1 = get_shape(slide, "Table 1")
    tbl5 = get_shape(slide, "Table 5")

    if not tbl1:
        print(f"  ✗ Table 1 not found in {src_file}")
        return

    raw_tbl = tbl1.table._tbl

    # 1. Get the original height of the single pricing row
    orig_rows = raw_tbl.findall(f"{{{NS}}}tr")
    orig_h    = int(orig_rows[PRICING_ROW_IDX].get("h", 0))

    # 2. Clone it twice so we have rows at PRICING_ROW_IDX, +1, +2
    clone_row_after(raw_tbl, PRICING_ROW_IDX)      # inserts at idx+1
    clone_row_after(raw_tbl, PRICING_ROW_IDX + 1)  # inserts at idx+2

    # 3. Set all 3 pricing rows to PRICING_ROW_H
    new_rows = raw_tbl.findall(f"{{{NS}}}tr")
    for i in range(PRICING_ROW_IDX, PRICING_ROW_IDX + 3):
        new_rows[i].set("h", str(PRICING_ROW_H))

    # 4. Move Table 5 down by the extra height we added
    extra = (3 * PRICING_ROW_H) - orig_h
    if tbl5 and extra > 0:
        tbl5.top = tbl5.top + extra

    prs.save(dst_file)
    orig_mb = os.path.getsize(src_file) / 1024 / 1024
    new_mb  = os.path.getsize(dst_file) / 1024 / 1024
    print(f"  ✓ {os.path.basename(dst_file)}  ({orig_mb:.1f}MB → {new_mb:.1f}MB)  extra={extra/914400*72:.1f}pt")

print("\nCreating tiered templates...")
make_tiered_template(
    os.path.join(BASE, "3_Mos_Pricing_Diligence_small.pptx"),
    os.path.join(BASE, "tiered_1term_small.pptx"),
)
make_tiered_template(
    os.path.join(BASE, "Datasite Proposal_6_9_options _small.pptx"),
    os.path.join(BASE, "tiered_2term_small.pptx"),
)
print("Done.\n")
