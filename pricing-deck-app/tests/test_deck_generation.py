#!/usr/bin/env python3
"""End-to-end smoke test for the Datasite Proposal Builder.

Exercises every pricing mode against a running server and validates that
the returned PowerPoint deck contains the form data that was submitted.

Two ways to run it:

    # with pytest (recommended for CI)
    pytest tests/ -v

    # standalone, no pytest required
    python3 tests/test_deck_generation.py

In either case the test auto-starts the Flask app on port 5050 if nothing
is already listening there, and tears it down on exit.

Environment overrides:
    PROPOSAL_BUILDER_URL   default http://localhost:5050
"""
from __future__ import annotations

import io
import os
import signal
import socket
import subprocess
import sys
import time
from pathlib import Path

import requests
from pptx import Presentation

BASE_DIR = Path(__file__).resolve().parent.parent
BASE_URL = os.environ.get("PROPOSAL_BUILDER_URL", "http://localhost:5050")


# ── Sample data ──────────────────────────────────────────────────────────────
# Canonical "Acme Corporation" scenario. The README references these same
# values in the manual walkthrough, so keep the two in sync.
SAMPLE_REP = {
    "signer_name":  "Shelle Martin",
    "signer_title": "Sales Director",
}
SAMPLE_DEAL = {
    "first_names": "Alex",
    "company":     "Acme Corporation",
    "industry":    "Technology",
    "prep":        "120",
    "pages1":      "10000",
    "pages2":      "20000",
}


def _payload(**extra) -> dict:
    p = {**SAMPLE_REP, **SAMPLE_DEAL}
    p.update(extra)
    return p


# ── Server lifecycle ─────────────────────────────────────────────────────────
def _port_open(host: str, port: int) -> bool:
    with socket.socket(socket.AF_INET, socket.SOCK_STREAM) as s:
        s.settimeout(0.5)
        try:
            s.connect((host, port))
            return True
        except OSError:
            return False


def _wait_healthy(url: str, timeout: float = 20.0) -> None:
    deadline = time.time() + timeout
    while time.time() < deadline:
        try:
            r = requests.get(url + "/healthz", timeout=2)
            if r.ok and r.json().get("ok"):
                return
        except requests.RequestException:
            pass
        time.sleep(0.3)
    raise RuntimeError(f"Server at {url} never became healthy within {timeout}s")


class _ServerHandle:
    """Context manager that starts `python3 app.py` if nothing is listening."""
    def __init__(self):
        self.proc: subprocess.Popen | None = None

    def __enter__(self):
        host_part = BASE_URL.split("://", 1)[-1]
        host, port = host_part.split(":")
        if _port_open(host, int(port)):
            return self
        self.proc = subprocess.Popen(
            [sys.executable, "app.py"],
            cwd=str(BASE_DIR),
            stdout=subprocess.DEVNULL,
            stderr=subprocess.DEVNULL,
        )
        _wait_healthy(BASE_URL)
        return self

    def __exit__(self, *exc):
        if self.proc is not None:
            try:
                self.proc.send_signal(signal.SIGINT)
                self.proc.wait(timeout=5)
            except Exception:
                self.proc.kill()


# ── PPTX inspection helpers ──────────────────────────────────────────────────
def _all_text(prs) -> str:
    lines = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                lines.append(shape.text_frame.text)
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        lines.append(cell.text)
    return "\n".join(lines)


def _find_table(prs, name: str):
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_table and shape.name == name:
                return shape.table
    return None


def _generate(payload: dict) -> Presentation:
    r = requests.post(BASE_URL + "/generate", data=payload, timeout=30)
    assert r.ok, f"POST /generate → HTTP {r.status_code}: {r.text[:300]}"
    ct = r.headers.get("Content-Type", "")
    assert ct.startswith("application/vnd.openxml"), f"Unexpected Content-Type: {ct}"
    assert len(r.content) > 100_000, f"Deck too small ({len(r.content)} bytes)"
    return Presentation(io.BytesIO(r.content))


def _assert_common_ingested(prs, payload: dict) -> None:
    """Fields every deck should contain regardless of pricing type."""
    text = _all_text(prs)
    assert payload["company"]     in text, "Company not ingested"
    assert payload["first_names"] in text, "First name not ingested"
    assert payload["signer_name"] in text, "Signer name not ingested"
    # Industry filtering: starting template has 7 industry slides; only the
    # chosen one should survive. The longest source template has 27 slides,
    # so a correctly filtered deck has at most 21.
    assert len(prs.slides) <= 21, (
        f"Industry filter did not fire (slide count={len(prs.slides)})"
    )
    # Retained industry slide should mention the industry keyword.
    keyword = payload["industry"].split()[0].lower()
    assert keyword in text.lower(), (
        f"Industry '{payload['industry']}' not present in retained slides"
    )


# ── Tests ────────────────────────────────────────────────────────────────────
# Each of these POSTs the form, asserts an .pptx came back, then inspects
# the deck to confirm the submitted data made it into the slides.

def test_healthz():
    r = requests.get(BASE_URL + "/healthz", timeout=5)
    assert r.ok, f"healthz returned {r.status_code}"
    assert r.json() == {"ok": True}


def test_industries_endpoint():
    r = requests.get(BASE_URL + "/industries", timeout=5)
    assert r.ok
    data = r.json()
    assert "industries" in data and isinstance(data["industries"], list)
    assert len(data["industries"]) == 7, (
        f"Expected 7 industries, got {len(data['industries'])}: {data['industries']}"
    )
    assert "Technology" in data["industries"]


def test_term1_flat_rate():
    payload = _payload(
        pricing_type="term1", num_terms="1",
        term1_months="3", term1_rate="0.37",
    )
    prs = _generate(payload)
    _assert_common_ingested(prs, payload)
    text = _all_text(prs)
    assert "$0.37/page" in text, "term1 rate not rendered"


def test_term2_two_flat_rates():
    payload = _payload(
        pricing_type="term2", num_terms="2",
        term2a_months="6", term2a_rate="0.59",
        term2b_months="9", term2b_rate="0.71",
    )
    prs = _generate(payload)
    _assert_common_ingested(prs, payload)
    text = _all_text(prs)
    assert "$0.59/page" in text, "term2a rate not rendered"
    assert "$0.71/page" in text, "term2b rate not rendered"


def test_warehouse_mtm():
    payload = _payload(
        pricing_type="mtm", num_terms="1",
        term1_months="3", term1_rate="0.15",
        mtm_rate="0.15",
    )
    prs = _generate(payload)
    _assert_common_ingested(prs, payload)
    text = _all_text(prs)
    assert "Warehouse" in text, "MTM not relabelled to 'Warehouse'"


def test_tiered_single_term():
    payload = _payload(
        pricing_type="tiered",
        num_tiers="3", tier_upper1="9999", tier_upper2="19999",
        tier1_rate="0.37", tier2_rate="0.32", tier3_rate="0.29",
    )
    prs = _generate(payload)
    _assert_common_ingested(prs, payload)
    text = _all_text(prs)
    for rate in ("$0.37/page", "$0.32/page", "$0.29/page"):
        assert rate in text, f"tiered rate {rate} missing"
    for band in ("Up to 9,999", "10,000–19,999", "20,000+"):
        assert band in text, f"tier band label {band!r} missing"


def test_tiered_two_terms():
    payload = _payload(
        pricing_type="tiered2",
        num_tiers="3", tier_upper1="9999", tier_upper2="19999",
        t2t_a_months="3", t2t_a_r1="0.40", t2t_a_r2="0.35", t2t_a_r3="0.30",
        t2t_b_months="6", t2t_b_r1="0.38", t2t_b_r2="0.33", t2t_b_r3="0.28",
    )
    prs = _generate(payload)
    _assert_common_ingested(prs, payload)
    text = _all_text(prs)
    for rate in ("$0.40/page", "$0.35/page", "$0.38/page", "$0.33/page"):
        assert rate in text, f"tiered2 rate {rate} missing"
    t = _find_table(prs, "Table 1")
    assert t is not None, "Pricing Table 1 not found"
    headers = [c.text for row in t.rows for c in row.cells]
    assert "3 Months" in headers and "6 Months" in headers


def test_tiered_three_terms():
    payload = _payload(
        pricing_type="tiered3",
        num_tiers="3", tier_upper1="9999", tier_upper2="19999",
        t3t_a_months="3", t3t_a_r1="0.59", t3t_a_r2="0.55", t3t_a_r3="0.50",
        t3t_b_months="6", t3t_b_r1="0.71", t3t_b_r2="0.65", t3t_b_r3="0.60",
        t3t_c_months="9", t3t_c_r1="0.85", t3t_c_r2="0.78", t3t_c_r3="0.70",
    )
    prs = _generate(payload)
    _assert_common_ingested(prs, payload)
    text = _all_text(prs)
    for rate in (
        "$0.59/page", "$0.55/page", "$0.50/page",
        "$0.71/page", "$0.65/page", "$0.60/page",
        "$0.85/page", "$0.78/page", "$0.70/page",
    ):
        assert rate in text, f"tiered3 rate {rate} missing"
    t = _find_table(prs, "Table 1")
    assert t is not None, "Pricing Table 1 not found"
    headers = [c.text for row in t.rows for c in row.cells]
    for h in ("3 Months", "6 Months", "9 Months"):
        assert h in headers, f"tiered3 table header {h!r} missing"
    t5 = _find_table(prs, "Table 5")
    assert t5 is not None, "Cost Table 5 not found"
    t5_headers = [c.text for row in t5.rows for c in row.cells]
    for h in ("3 Month", "6 Month", "9 Month"):
        assert h in t5_headers, f"tiered3 cost header {h!r} missing"


def test_custom_tier_boundaries_2tier():
    """Custom boundaries must propagate into band labels (2-tier mode)."""
    payload = _payload(
        pricing_type="tiered2", num_tiers="2",
        tier_upper1="5000",
        t2t_a_months="3", t2t_a_r1="0.40", t2t_a_r2="0.35", t2t_a_r3="0.0",
        t2t_b_months="6", t2t_b_r1="0.38", t2t_b_r2="0.33", t2t_b_r3="0.0",
    )
    prs = _generate(payload)
    text = _all_text(prs)
    assert "Up to 5,000" in text, "Custom tier_upper1 not reflected in band label"
    assert "5,001+"    in text, "Second-band label not computed from custom bound"


def test_bad_input_returns_400():
    """Malformed numeric input returns a friendly 400, not a 500."""
    payload = _payload(
        pricing_type="term1", num_terms="1",
        term1_months="3", term1_rate="abc",   # ← invalid
    )
    r = requests.post(BASE_URL + "/generate", data=payload, timeout=10)
    assert r.status_code == 400, f"Expected 400, got {r.status_code}"
    assert "Check your inputs" in r.text


# ── Standalone runner (no pytest required) ───────────────────────────────────
TESTS = [
    test_healthz,
    test_industries_endpoint,
    test_term1_flat_rate,
    test_term2_two_flat_rates,
    test_warehouse_mtm,
    test_tiered_single_term,
    test_tiered_two_terms,
    test_tiered_three_terms,
    test_custom_tier_boundaries_2tier,
    test_bad_input_returns_400,
]


def _main() -> int:
    with _ServerHandle():
        failures: list[tuple[str, str]] = []
        for t in TESTS:
            try:
                t()
            except AssertionError as e:
                failures.append((t.__name__, str(e) or "assertion failed"))
                print(f"FAIL  {t.__name__}\n      {e}")
            except Exception as e:
                failures.append((t.__name__, repr(e)))
                print(f"ERROR {t.__name__}\n      {e!r}")
            else:
                print(f"PASS  {t.__name__}")
        print()
        print(f"{len(TESTS) - len(failures)} passed, {len(failures)} failed")
        return 1 if failures else 0


if __name__ == "__main__":
    sys.exit(_main())
